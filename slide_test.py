from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def clear_paragraph(paragraph):
    for run in paragraph.runs:
        run.text = ''


def update_title_slide(prs, new_title, new_subtitle):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "$title" in shape.text:
                    shape.text = new_title
                elif "$titlesubtitle" in shape.text:
                    shape.text = new_subtitle


# def update_bullet_point_slide(prs, new_title, new_bulletpoint1, new_bulletpoint2):
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 if "$bulletpointslide" in shape.text:
#                     shape.text = new_title
#                 if "$bulletpoint1" in shape.text:
#                     shape.text = new_bulletpoint1
#                 elif "$bulletpoint2" in shape.text:
#                     shape.text = new_bulletpoint2

def update_bullet_point_slide(prs, new_title, new_bulletpoints):
    points = new_bulletpoints.split(".")
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if "$bulletpointslide" in paragraph.text:
                        paragraph.text = new_title
                    elif "$bulletpoint1" in paragraph.text:
                        paragraph.shape = slide.shapes.placeholders[1]
                        paragraph.text = "d"

                        # clear_paragraph(paragraph)

                        for point in points:
                            print(point)
                            p = shape.text_frame.add_paragraph()
                            p.shape = slide.shapes.placeholders[1]
                            p.text = "whats good!"
                            # p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

                    # elif "$bulletpoint2" in paragraph.text:
                    #     # Clear existing text but preserve formatting
                    #     clear_paragraph(paragraph)
                    #
                    #     # Add new bullet points
                    #     for point in points:
                    #         print(point)
                    #         p = shape.text_frame.add_paragraph()
                    #         p.text = point
                    #         p.level = 1
                    #         p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

# Add more functions for other slide types...

# Load the presentation
presentation_path = 'themes/layoutslides.pptx'
prs = Presentation(presentation_path)

# Update slides
update_title_slide(prs, "New Title", "New Subtitle")
# update_bullet_point_slide(prs, "New Bullet Point Slide Title", "New Bullet Point 1", "New Bullet Point 2")
update_bullet_point_slide(prs, "title woohoo", "asdjopasjdpasdjaiposj. dsajdajsdjaskldj.")
# Call other functions for different slides here...

# Save the updated presentation to a new file
prs.save('new_presentation.pptx')