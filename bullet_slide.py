from pptx import Presentation


class BulletSlide:

    def __init__(self, prs, title, bulletpoint1, subbulletpoint1, bulletpoint2, subbulletpoint2, bulletpoint3, subbulletpoint3):
        self.prs = prs
        self.title = title
        self.bulletpoint1 = bulletpoint1
        self.subbulletpoint1 = subbulletpoint1
        self.bulletpoint2 = bulletpoint2
        self.subbulletpoint2 = subbulletpoint2
        self.bulletpoint3 = bulletpoint3
        self.subbulletpoint3 = subbulletpoint3
        self.bullet_slide_layout = prs.slide_layouts[1]


    def createBulletSlide(self):
        slide = self.prs.slides.add_slide(self.bullet_slide_layout)

        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = self.title

        tf = body_shape.text_frame
        tf.text = self.bulletpoint1

        p = tf.add_paragraph()
        p.text = self.subbulletpoint1
        p.level = 1

        # p = tf.add_paragraph()
        # p.text = self.subbulletpoint3
        # p.level = 0



        return slide
