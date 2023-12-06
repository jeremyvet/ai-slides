import copy
from pptx import Presentation


def compare_text(text_frame, text):
    return text_frame.paragraphs[0].text == text


def clear_slide(slide):
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)


def set_text(text_frame, new_text):
    p = text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = new_text
        for run in p.runs[1:]:
            p._element.remove(run._element)
    else:
        p.add_run().text = new_text



class SlideshowGenerator:
    def __init__(self):
        self.pres = Presentation()
        self.save_file_name = ""
        self.save_callbacks = []
    
    def copy_slide(self, source, idx):
        ext_slide = source.slides[idx]

        self.pres.slide_width = source.slide_width
        self.pres.slide_height = source.slide_height

        SLD_LAYOUT = 5
        slide_layout = self.pres.slide_layouts[SLD_LAYOUT]

        curr_slide = self.pres.slides.add_slide(slide_layout)

        clear_slide(curr_slide)

        for shp in ext_slide.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        if ext_slide.background.fill.type == 1:  # Solid fill
            fill = curr_slide.background.fill
            fill.solid()
            fill.fore_color.rgb = ext_slide.background.fill.fore_color.rgb

        return source

    def create_title(self, theme: str, title: str, subtitle: str):
        pres = Presentation("./themes/" + theme + '.pptx')

        for shape in pres.slides[0].shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame

                if compare_text(text_frame, "$text.title"):
                    set_text(text_frame, title)
                elif compare_text(text_frame, "$text.subtitle"):
                    set_text(text_frame, subtitle)

        self.copy_slide(pres, 0)

    def create_content_slide(self, theme: str, title: str, content: str):
        pres = Presentation("./themes/" + theme + '.pptx')

        for shape in pres.slides[1].shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame

                if compare_text(text_frame, "$text.title"):
                    set_text(text_frame, title)
                elif compare_text(text_frame, "$text.content"):
                    set_text(text_frame, content)

        self.copy_slide(pres, 1)

    def save(self):
        self.pres.save("./presentations/" + self.save_file_name + ".pptx")

        for callback in self.save_callbacks:
            callback()
