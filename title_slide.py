from pptx import Presentation


class TitleSlide:

    def __init__(self, prs, title, subtitle):
        self.prs = prs
        self.title = title
        self.subtitle = subtitle
        self.title_slide_layout = prs.slide_layouts[0]

    def createTitleSlide(self):
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        slide.shapes.title.text = self.title
        slide.placeholders[1].text = self.subtitle
        return slide
