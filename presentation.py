from pptx import Presentation
from title_slide import TitleSlide
from bullet_slide import BulletSlide

def add_slide(prs, layout, title):
    """add slides with inputted values"""
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    return slide

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]


# slide = add_slide(prs, title_slide_layout, "Title Slide")
slide = TitleSlide(prs, "Title Test", "test Tickler").createTitleSlide()
slide2 = add_slide(prs, bullet_slide_layout, "Bullet Point Slide")
slide3 = BulletSlide(prs, "Bullet Test", "bullet point 1", "sub bullet point 1", "bullet point 2", "e", "e", "e")

prs.save('test.pptx')
#
# title_slide_layout = prs.slide_layouts[0]
#
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
#
# title.text = "testing pptx"
# subtitle.text = "trash documentation"
#
# bullet_slide_layout = prs.slide_layouts[1]
# slide2 = prs.slides.add_slide(bullet_slide_layout)
# shapes = slide.shapes
#
# title_shape = shapes.title
# body_shape = shapes.placeholders[1]
#
# title_shape.text = "bullet text example"
#
# tf = body_shape.text_frame
# tf.text = 'suck my massive red nuts'
#
# p = tf.add_paragraph()
# p.text = 'bullet point 2'
# p.level = 1
#
# p = tf.add_paragraph()
# p.text = 'bullet point 3'
# p.level = 2
#
# prs.save('test.pptx')
